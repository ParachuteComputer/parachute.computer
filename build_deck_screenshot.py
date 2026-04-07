"""Build NVC pitch deck PPTX from screenshots of the web deck.

Generates a Keynote-compatible PPTX where each slide is a full-bleed
screenshot of the corresponding web slide. Speaker notes are pulled
from the .speaker-notes elements in the source njk and embedded in
the PPTX notes pane.

Why screenshots? Keynote on macOS has long-standing compatibility issues
with python-pptx-generated PPTX files. Screenshot-based PPTX files have
no such issues — they're just images, which every tool reads cleanly.

Usage:
  python3 build_deck_screenshot.py
"""

import http.server
import os
import re
import shutil
import socketserver
import subprocess
import tempfile
import threading
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).parent
SITE_DIR = ROOT / "_site"
DECK_NJK = ROOT / "nvc-deck.njk"
OUT_PPTX = ROOT / "nvc-pitch-deck.pptx"
SERVED_PPTX = ROOT / "nvc" / "Parachute.pptx"

PORT = 8765
DECK_URL = f"http://localhost:{PORT}/nvc/deck/"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Capture viewport (will be 2x for retina-quality)
VIEWPORT_W = 1280
VIEWPORT_H = 720
DEVICE_SCALE = 2


def build_eleventy():
    print("Building eleventy site...")
    subprocess.run(["npx", "@11ty/eleventy"], cwd=ROOT, check=True)


class _Handler(http.server.SimpleHTTPRequestHandler):
    def log_message(self, format, *args):
        pass  # silence


def _start_server():
    os.chdir(SITE_DIR)
    handler = _Handler
    server = socketserver.TCPServer(("localhost", PORT), handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    return server


def screenshot_slides(output_dir: str) -> list[str]:
    server = _start_server()
    paths = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            ctx = browser.new_context(
                viewport={"width": VIEWPORT_W, "height": VIEWPORT_H},
                device_scale_factor=DEVICE_SCALE,
            )
            page = ctx.new_page()
            page.goto(DECK_URL, wait_until="networkidle")

            # Hide download bar, notes toggle, and speaker notes
            page.add_style_tag(content="""
                div[style*='position: fixed'] { display: none !important; }
                .notes-toggle { display: none !important; }
                .speaker-notes { display: none !important; }
                html { scroll-snap-type: none !important; }
            """)

            page.wait_for_timeout(800)

            slides = page.query_selector_all("section.slide")
            print(f"Found {len(slides)} slides")

            for i, slide in enumerate(slides):
                # Force the slide to fill exactly the viewport for clean capture
                page.evaluate("""(el) => {
                    el.style.minHeight = '100vh';
                    el.style.height = '100vh';
                    el.style.maxHeight = '100vh';
                    el.style.overflow = 'hidden';
                    el.style.paddingTop = '4rem';
                    el.style.paddingBottom = '4rem';
                }""", slide)

                slide.scroll_into_view_if_needed()
                page.wait_for_timeout(300)

                box = slide.bounding_box()
                path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
                page.screenshot(
                    path=path,
                    clip={
                        "x": 0,
                        "y": box["y"],
                        "width": VIEWPORT_W,
                        "height": VIEWPORT_H,
                    },
                )
                paths.append(path)
                print(f"  Captured slide {i+1}")

            browser.close()
    finally:
        server.shutdown()
        os.chdir(ROOT)

    return paths


def extract_speaker_notes() -> list[str]:
    """Pull speaker notes from the source njk file."""
    with open(DECK_NJK, "r") as f:
        content = f.read()

    pattern = (
        r'<div class="speaker-notes">\s*'
        r'<div class="notes-label">Speaker Notes</div>\s*'
        r'(.*?)\s*'
        r'</div>'
    )
    matches = re.findall(pattern, content, re.DOTALL)

    notes = []
    for match in matches:
        text = match.strip()
        # Strip HTML tags
        text = re.sub(r"<[^>]+>", "", text)
        # Decode common entities
        text = (text
                .replace("&mdash;", "—")
                .replace("&ndash;", "–")
                .replace("&middot;", "·")
                .replace("&rsquo;", "'")
                .replace("&lsquo;", "'")
                .replace("&ldquo;", '"')
                .replace("&rdquo;", '"')
                .replace("&amp;", "&")
                .replace("&nbsp;", " "))
        # Collapse whitespace
        text = re.sub(r"\s+", " ", text).strip()
        notes.append(text)

    return notes


def build_pptx(image_paths: list[str], speaker_notes: list[str]):
    print("Building PPTX...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(image_paths):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, SLIDE_W, SLIDE_H)

        if i < len(speaker_notes):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes[i]

    prs.save(OUT_PPTX)
    SERVED_PPTX.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(OUT_PPTX, SERVED_PPTX)
    print(f"Saved {OUT_PPTX}")
    print(f"Copied to {SERVED_PPTX}")
    size_mb = os.path.getsize(OUT_PPTX) / 1024 / 1024
    print(f"Size: {size_mb:.1f} MB")


def main():
    build_eleventy()
    speaker_notes = extract_speaker_notes()
    print(f"Extracted {len(speaker_notes)} speaker notes")

    with tempfile.TemporaryDirectory() as tmpdir:
        images = screenshot_slides(tmpdir)
        build_pptx(images, speaker_notes)


if __name__ == "__main__":
    main()
