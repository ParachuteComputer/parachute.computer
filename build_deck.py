"""Build NVC pitch deck for Parachute.
10 slides, image-forward, with speaker notes."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Brand colors ──
BG       = RGBColor(0xFA, 0xF8, 0xF4)
BG_SOFT  = RGBColor(0xF3, 0xF0, 0xEA)
FG       = RGBColor(0x2C, 0x2A, 0x26)
FG_MUTED = RGBColor(0x6B, 0x68, 0x60)
FG_DIM   = RGBColor(0x9A, 0x96, 0x90)
ACCENT   = RGBColor(0x4A, 0x7C, 0x59)
ACCENT_SOFT = RGBColor(0xE8, 0xF1, 0xEA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0xE4, 0xE0, 0xD8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SERIF = "Georgia"
FONT_SANS = "Calibri"
FONT_MONO = "Courier New"

IMAGES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "nvc", "images")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]


# ─────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=FG, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_SANS, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return tf


def section_label(slide, text, top=Inches(0.9), left=Inches(0.8)):
    add_text_box(slide, left, top, Inches(5), Inches(0.4),
                 text.upper(), font_size=11, color=ACCENT, bold=True,
                 font_name=FONT_SANS)


def slide_headline(slide, text, top=Inches(1.4), size=42, left=Inches(0.8), width=Inches(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = FG
        p.font.bold = False
        p.font.name = FONT_SERIF
        p.line_spacing = 1.05
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return tf


def slide_sub(slide, text, top, left=Inches(0.8), width=Inches(8), size=15, color=FG_MUTED):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_SANS
        p.line_spacing = 1.4
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return tf


def add_image(slide, filename, left, top, width=None, height=None):
    path = os.path.join(IMAGES_DIR, filename)
    if not os.path.exists(path):
        print(f"  WARNING: missing {filename}")
        return None
    if width and height:
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(path, left, top, width=width)
    if height:
        return slide.shapes.add_picture(path, left, top, height=height)
    return slide.shapes.add_picture(path, left, top)


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def add_footer(slide, text="parachute.computer"):
    add_text_box(slide, Inches(0.8), Inches(7.05), Inches(5), Inches(0.3),
                 text, font_size=9, color=FG_DIM, font_name=FONT_SANS)


def add_brand_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def new_slide():
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s)
    add_brand_bar(s)
    add_footer(s)
    return s


# ═══════════════════════════════════════════════
# SLIDE 1 — ASPIRATION
# ═══════════════════════════════════════════════
s = new_slide()

add_text_box(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
             "Parachute", font_size=28, color=FG, font_name=FONT_SERIF)

slide_headline(s, "Capture. Connect. Compound.", top=Inches(2.2), size=48, width=Inches(8))

slide_sub(s, "A knowledge layer that makes any AI truly yours.",
          top=Inches(4.4), size=18)

add_image(s, "slide1.png", Inches(7.5), Inches(1.0), width=Inches(5.0))

add_speaker_notes(s, (
    "A thought on a walk. A note after a meeting. An idea at 2am. "
    "Right now, those are scattered — lost in apps that don't talk to each other "
    "and invisible to the AI tools we rely on. Parachute captures your thinking "
    "wherever you are, connects it into a structured knowledge graph, and "
    "compounds it over time into something that makes every AI you use more "
    "personal, more capable, more yours. I'm wearing one of our capture devices "
    "right now. This is Parachute."
))


# ═══════════════════════════════════════════════
# SLIDE 2 — THE RACE
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Landscape")
slide_headline(s, "Everyone is building\nAI agents.", top=Inches(1.4), size=44)

slide_sub(s, "OpenClaw. Claude Cowork. Perplexity Computer. Manus.\n"
             "We were too. Features get cloned in weeks.\n"
             "So we stopped trying to win that race.",
          top=Inches(4.0), size=16)

add_image(s, "slide2.png", Inches(7.2), Inches(1.4), width=Inches(5.5))

add_speaker_notes(s, (
    "OpenClaw. Claude Cowork. Perplexity Computer. Manus. There are dozens of "
    "companies racing to build the best AI agent. Over 100 million people already "
    "pay $20+/month for AI. We started here too — we were building an agentic "
    "computing platform. But we watched features get cloned in weeks. Every new "
    "capability one platform shipped, three others had within a month. It's a "
    "race to the bottom, and the big labs have billions of dollars to outspend "
    "everyone. So we stopped trying to win that race. We asked: what's actually "
    "missing? What do all of these agents need that none of them are building?"
))


# ═══════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Problem")
slide_headline(s, "There's no\nknowledge layer.", top=Inches(1.4), size=44)

slide_sub(s, "Every AI stores your context as one flat text file.\n"
             "And to give it context, you have to talk to it.\n"
             "There's no way to just think and have it become context.",
          top=Inches(4.2), size=15)

add_image(s, "slide3.jpeg", Inches(7.2), Inches(1.6), width=Inches(5.5))

add_speaker_notes(s, (
    "Every AI platform — ChatGPT, Claude, Gemini — stores your personal context "
    "the same way: one big text file. It remembers your name, maybe a few "
    "preferences. It doesn't know six months of your thinking. And to give it "
    "context, you have to talk to it directly. There's no way to just think for "
    "yourself — on a walk, in the car — and have that become context your AI "
    "can use. You might ask, 'why not just dictate into ChatGPT?' Because then "
    "your thoughts are trapped inside one platform, unstructured, unsearchable, "
    "and gone the moment you switch tools. The knowledge layer between you and "
    "your AI doesn't exist. We're building it."
))


# ═══════════════════════════════════════════════
# SLIDE 4 — THE PRODUCT
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Product")
slide_headline(s, "Parachute Daily.", top=Inches(1.4), size=48)

slide_sub(s, "Voice-first journal. Talk or type. Phone, web, or pendant.\n"
             "Open source · Local-first · Your data is yours.",
          top=Inches(3.5), size=16)

add_image(s, "slide4.png", Inches(7.0), Inches(1.4), width=Inches(5.7))

add_speaker_notes(s, (
    "Parachute Daily is a voice-first journal. Talk or type — on your phone, "
    "on the web, or through a wearable pendant. Your words get transcribed and "
    "structured into a personal knowledge graph: notes, tags, and links. Simple "
    "for humans, native for AI. Most people will start with the app on their "
    "phone. The pendant is an accessory for people who want frictionless voice "
    "capture — press a button, talk on a walk, your thoughts are structured by "
    "the time you're home. Working prototype, on stage today. On privacy: this "
    "is a Public Benefit Corporation. The code is open source. Your data is "
    "local-first — it lives on your device by default. The cloud is opt-in, "
    "and if you choose it, we use end-to-end encryption. You can export "
    "everything, self-host, or leave at any time. We designed this for trust "
    "from day one, not as an afterthought."
))


# ═══════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "How It Works")
slide_headline(s, "Spin up a vault.\nGive any AI a link.", top=Inches(1.4), size=42)

slide_sub(s, "Built on MCP — the open standard for connecting AI to tools.\n"
             "Your data stays yours. Works with every AI, not instead of one.",
          top=Inches(4.2), size=14)

add_image(s, "slide5.png", Inches(2.5), Inches(5.2), width=Inches(8.0))

# If we have room — make it bigger / center it
# Actually, place the diagram center-bottom
# Let me re-think the layout: top has headline+sub, bottom has the diagram

add_speaker_notes(s, (
    "Your data lives in a Parachute Vault — a personal knowledge graph organized "
    "around notes, tags, and links. It's accessible to any AI via MCP, the open "
    "standard Anthropic created for connecting AI to tools. You spin up a vault "
    "in seconds and give any AI a link. That's the whole setup. Whatever AI you "
    "already use — Claude, ChatGPT, Gemini — Parachute makes it better. You "
    "don't switch anything. You add Parachute. This is how we stay independent "
    "from any one LLM. We're not tied to a single model. We speak an open "
    "protocol that any AI can use. As users settle into different AI camps, "
    "we're the layer that works across all of them. That's why we're not in "
    "the race to the bottom on AI features — we're building the infrastructure "
    "underneath it."
))


# ═══════════════════════════════════════════════
# SLIDE 6 — COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Competitive Landscape")
slide_headline(s, "Notion traps you.\nObsidian loses you.", top=Inches(1.4), size=42, width=Inches(8))

slide_sub(s, "Notion is where 100M people store their thinking — locked inside\n"
             "one platform, accessible only to their AI. Obsidian is where\n"
             "power users pair structured notes with AI — but it takes\n"
             "markdown, plugins, and a command line.\n"
             "\n"
             "Parachute is human-native and agent-native. The road between.",
          top=Inches(3.6), size=13, width=Inches(7))

add_image(s, "slide6.png", Inches(7.5), Inches(2.0), width=Inches(5.3))

add_speaker_notes(s, (
    "Two big competitors. Notion has 100 million users. It's where teams and "
    "individuals are increasingly storing their knowledge, and people are using "
    "AI to create Notion content all the time. But Notion is a walled garden "
    "— your data is locked inside their platform, and the AI that runs on it "
    "is theirs. The moment you want to use Claude or ChatGPT or whatever's "
    "next, you're copy-pasting things out. Obsidian is the most interesting "
    "thing happening at the power-user end. Bootstrapped, profitable, beloved "
    "— and a growing community is pairing it with Claude Code to do incredible "
    "things with structured knowledge plus AI. But it requires markdown files, "
    "plugin configuration, command-line tools. 99% of people will never do "
    "that. So Notion traps you and Obsidian loses you. Parachute is the third "
    "way. Human-native — voice-first, simple, no setup. Agent-native — built "
    "on MCP, works with any AI you choose. Your knowledge layer comes with "
    "you, regardless of which AI you use today or tomorrow. We're not picking "
    "a fight with Notion or Obsidian. We're the open road between them."
))


# ═══════════════════════════════════════════════
# SLIDE 7 — WHO USES THIS
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Who Uses This")
slide_headline(s, "Who needs this first?", top=Inches(1.4), size=44, width=Inches(12))

add_image(s, "slide7.png", Inches(1.5), Inches(3.2), width=Inches(10.3))

add_speaker_notes(s, (
    "Our beachhead customer is the AI journaler — someone who already uses "
    "Claude or ChatGPT daily and journals or takes voice notes, but is "
    "frustrated that their AI doesn't know what they were thinking about "
    "yesterday. That's our first 500 paid users. We reach them through Learn "
    "Vibe Build, our AI learning cohorts, and through the Claude and ChatGPT "
    "power-user communities. Next is the professional — voice-notes after "
    "meetings, months of institutional context, searchable by any tool, never "
    "trapped in one app. Then enterprise: a departing team member records days "
    "of knowledge transfer. Their expertise lives on in a structured, queryable "
    "vault. Think about retiring baby boomers with decades of institutional "
    "knowledge — companies will pay for a smooth way to capture that. Each "
    "stage expands the market, but we start narrow and specific."
))


# ═══════════════════════════════════════════════
# SLIDE 8 — BUSINESS MODEL + MARKET
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Business Model")
slide_headline(s, "The $2–10 add-on.", top=Inches(1.4), size=44, width=Inches(7))

# Tier table on the left
tier_left = Inches(0.8)
tier_top = Inches(3.0)
tier_w = Inches(6.5)
tier_row_h = Inches(0.55)

tiers = [
    ("Free", "Offline journal", False),
    ("$2/mo", "Cloud sync + MCP", False),
    ("$5/mo", "Cloud transcription", False),
    ("$10/mo", "AI reflections + synthesis", True),
    ("$100", "Pendant", False),
]

# Outer container
tier_outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                tier_left, tier_top, tier_w,
                                tier_row_h * len(tiers) + Inches(0.1))
tier_outer.fill.solid()
tier_outer.fill.fore_color.rgb = WHITE
tier_outer.line.color.rgb = BORDER
tier_outer.line.width = Pt(0.75)

for i, (price, desc, highlight) in enumerate(tiers):
    row_top = tier_top + Inches(0.05) + (tier_row_h * i)

    if highlight:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                tier_left + Inches(0.05), row_top,
                                tier_w - Inches(0.1), tier_row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = ACCENT_SOFT
        bg.line.fill.background()

    add_text_box(s, tier_left + Inches(0.4), row_top + Inches(0.12),
                 Inches(1.8), Inches(0.4),
                 price, font_size=18, color=ACCENT, font_name=FONT_SERIF)
    add_text_box(s, tier_left + Inches(2.2), row_top + Inches(0.18),
                 Inches(4), Inches(0.4),
                 desc, font_size=13, color=FG_MUTED, bold=True)

# Right side: market stat + projections
right_left = Inches(8.0)

# Big stat
add_text_box(s, right_left, Inches(2.0), Inches(5), Inches(1.2),
             "100M+", font_size=52, color=ACCENT, font_name=FONT_SERIF)
add_text_box(s, right_left, Inches(3.1), Inches(5), Inches(0.8),
             "people pay $20+/mo for AI already",
             font_size=12, color=FG_MUTED)

# Projections table
proj_top = Inches(4.2)
proj_left = right_left
proj_w = Inches(4.8)

proj_outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                proj_left, proj_top, proj_w, Inches(2.2))
proj_outer.fill.solid()
proj_outer.fill.fore_color.rgb = WHITE
proj_outer.line.color.rgb = BORDER
proj_outer.line.width = Pt(0.75)

# Header
header_top = proj_top + Inches(0.1)
add_text_box(s, proj_left + Inches(0.2), header_top, Inches(1.5), Inches(0.3),
             "", font_size=9, color=FG_DIM, bold=True)
add_text_box(s, proj_left + Inches(1.7), header_top, Inches(1), Inches(0.3),
             "2026", font_size=9, color=FG_DIM, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, proj_left + Inches(2.7), header_top, Inches(1), Inches(0.3),
             "2027", font_size=9, color=FG_DIM, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, proj_left + Inches(3.7), header_top, Inches(1), Inches(0.3),
             "2028", font_size=9, color=FG_DIM, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ("Paid users", "500", "8,000", "50,000"),
    ("ARR", "~$30K", "~$480K", "~$3M"),
    ("Opex", "~$160K", "~$500K", "~$1.2M"),
]

for i, (label, c1, c2, c3) in enumerate(rows):
    row_top = proj_top + Inches(0.5) + Inches(0.45 * i)
    add_text_box(s, proj_left + Inches(0.2), row_top, Inches(1.5), Inches(0.3),
                 label, font_size=11, color=FG, bold=True)
    is_arr = label == "ARR"
    color = ACCENT if is_arr else FG_MUTED
    bold = is_arr
    add_text_box(s, proj_left + Inches(1.7), row_top, Inches(1), Inches(0.3),
                 c1, font_size=11, color=color, bold=bold, alignment=PP_ALIGN.CENTER)
    add_text_box(s, proj_left + Inches(2.7), row_top, Inches(1), Inches(0.3),
                 c2, font_size=11, color=color, bold=bold, alignment=PP_ALIGN.CENTER)
    add_text_box(s, proj_left + Inches(3.7), row_top, Inches(1), Inches(0.3),
                 c3, font_size=11, color=color, bold=bold, alignment=PP_ALIGN.CENTER)

add_speaker_notes(s, (
    "We're not competing with anyone's AI subscription. Over 100 million people "
    "already pay $20+/mo for AI. We're the $2–10 add-on that makes it "
    "dramatically better. Free tier has zero hosting cost — it's a fully "
    "offline journal. Each paid tier adds cloud features. The pendant is a "
    "one-time hardware purchase. Lower COGS than agentic platforms because "
    "we're not running inference. The path from 500 to 50,000: our first 500 "
    "come from three channels — Learn Vibe Build AI cohorts, the Claude and "
    "ChatGPT power-user communities, and our Boulder co-working space as a "
    "local testbed. Those 500 users generate the case studies and word-of-mouth "
    "that fuel organic growth. By year two, we add the open source developer "
    "community building MCP integrations, and early professional users. Year "
    "three, enterprise pilots for knowledge capture. Each channel compounds. "
    "We're not relying on paid acquisition. CAC stays low because the product "
    "sells itself once someone sees what their AI can do with six months of "
    "structured context."
))


# ═══════════════════════════════════════════════
# SLIDE 9 — TEAM & TRACTION
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Team & Traction")
slide_headline(s, "Self-funded.\nBuilt from scratch.", top=Inches(1.4), size=44)

slide_sub(s, "Working app · MCP server · Pendant prototype · Beta launching this month",
          top=Inches(4.0), size=13)

team = [
    ("Aaron Neyer", "Co-founder · Product & architecture"),
    ("Jon Bo", "Co-founder · Parachute Daily lead"),
    ("Lucian Hymer", "AI & infrastructure"),
    ("Marvin Melzer", "Pendant hardware"),
    ("Neil Yarnal", "Brand & design"),
]

team_top = Inches(4.6)
for i, (name, role) in enumerate(team):
    row_top = team_top + Inches(0.4 * i)
    add_text_box(s, Inches(0.8), row_top, Inches(3), Inches(0.35),
                 name, font_size=14, color=FG, font_name=FONT_SERIF)
    add_text_box(s, Inches(3.5), row_top + Inches(0.05), Inches(7), Inches(0.35),
                 "— " + role, font_size=11, color=FG_MUTED)

add_speaker_notes(s, (
    "Everything you've seen is working today. App with local voice transcription "
    "and graph storage. MCP server. Pendant prototype — I'm wearing it right "
    "now. Daily beta launching this month. PBC incorporated in Colorado. Me: "
    "finishing my MS at CU ATLAS in Creative Technology & Design, also hold an "
    "MA in Ecopsychology. Ex-Google. 10+ years full stack. Jon Bo is my "
    "co-founder and leads Parachute Daily — he's been a 3x founding engineer, "
    "and he's also an avid note-taker, journaler, and writer with an incredibly "
    "strong product sense and a gift for simple, stable architecture. He's the "
    "ideal person to build this. Lucian Hymer leads AI and infrastructure. "
    "Marvin Melzer built the pendant prototype. Neil Yarnal handles brand and "
    "design. We've built everything you see self-funded."
))


# ═══════════════════════════════════════════════
# SLIDE 10 — THE ASK
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Ask")
slide_headline(s, "$300K to launch\nthe knowledge layer.", top=Inches(1.4), size=44)

# Ask box
ask_left = Inches(0.8)
ask_top = Inches(4.0)
ask_w = Inches(6)
ask_h = Inches(2.0)

ask_outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               ask_left, ask_top, ask_w, ask_h)
ask_outer.fill.solid()
ask_outer.fill.fore_color.rgb = WHITE
ask_outer.line.color.rgb = BORDER
ask_outer.line.width = Pt(0.75)

ask_items = [
    ("Instrument", "SAFE"),
    ("Use of funds", "Core team through 2026"),
    ("Goal", "Launch by June · revenue immediately"),
]

for i, (label, val) in enumerate(ask_items):
    row_top = ask_top + Inches(0.25) + Inches(0.55 * i)
    add_text_box(s, ask_left + Inches(0.4), row_top, Inches(2), Inches(0.3),
                 label, font_size=11, color=FG_MUTED)
    add_text_box(s, ask_left + Inches(2.5), row_top, Inches(3.5), Inches(0.3),
                 val, font_size=14, color=FG, font_name=FONT_SERIF)

# Closing quote
add_text_box(s, Inches(7.5), Inches(3.2), Inches(5.2), Inches(2),
             "A knowledge layer for humans and AI",
             font_size=22, color=FG, font_name=FONT_SERIF, line_spacing=1.3)
add_text_box(s, Inches(7.5), Inches(3.65), Inches(5.2), Inches(2),
             "to think and remember together.",
             font_size=22, color=FG, font_name=FONT_SERIF, line_spacing=1.3)

# Contact
add_text_box(s, Inches(7.5), Inches(5.2), Inches(5.5), Inches(0.4),
             "Aaron Gabriel Neyer  ·  aaron@parachute.computer",
             font_size=10, color=FG_DIM)
add_text_box(s, Inches(7.5), Inches(5.45), Inches(5.5), Inches(0.4),
             "(513) 593-1721  ·  Boulder, CO",
             font_size=10, color=FG_DIM)
add_text_box(s, Inches(7.5), Inches(5.75), Inches(5.5), Inches(0.4),
             "parachute.computer  ·  github.com/ParachuteComputer",
             font_size=10, color=ACCENT)

add_speaker_notes(s, (
    "We're raising $300K on a SAFE to get the core team full-time through 2026. "
    "Production launch by June with revenue immediately. A knowledge layer for "
    "humans and AI to think and remember together."
))


# ─────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "nvc-pitch-deck.pptx")
prs.save(out_path)
print(f"Saved to {out_path}")
